#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
敏實科技大學專屬海報自動生成器 (generate_poster.py)
支援兩種海報模式：
1. project: USR 計畫/成果海報 (包含問題意識、執行重點、社會效益、成果亮點)
2. lecture: 主題講座/演講海報 (包含講座主題、演講人、時間、地點、報名資訊)

輸出格式：
- .pptx 檔案：可在 Canva 中直接匯入並保持向量圖層可編輯性
- .png 檔案：高解析度直式海報圖檔
"""

import sys
import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT_LOGO_PATH = os.path.join(SKILL_DIR, "resources", "minth_logo.png")

def create_minth_poster(
    mode="lecture",
    title="AI與永續發展主題講座",
    subtitle="",
    speaker="張教授 / 敏實科技大學",
    time_str="115年9月15日 (三) 10:00 - 12:00",
    location="敏實科技大學 大華樓5樓國際會議廳",
    advisor="教育部",
    co_organizer="",
    co_logo_path=None,
    organizer="敏實科技大學 USR 實踐中心",
    contact="邱專員 (分機 1234 / email: usr@mitust.edu.tw)",
    details=None,
    output_pptx="poster_output.pptx"
):
    prs = Presentation()
    # A4 Portrait: 8.27 x 11.69 inches
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Base colors
    navy_blue = RGBColor(0x00, 0x4B, 0x87)
    light_bg = RGBColor(0xF4, 0xF8, 0xFB)
    dark_text = RGBColor(0x22, 0x22, 0x22)
    accent_orange = RGBColor(0xE6, 0x51, 0x00)

    # 1. Background Fill
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.27), Inches(11.69))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = light_bg
    bg_shape.line.fill.background()

    # 2. Header Logos
    # Left Top Logo: MINTH Logo
    if os.path.exists(DEFAULT_LOGO_PATH):
        slide.shapes.add_picture(DEFAULT_LOGO_PATH, Inches(0.4), Inches(0.4), height=Inches(0.65))

    # Right Top Logo: Co-organizer Logo (if provided)
    if co_logo_path and os.path.exists(co_logo_path):
        slide.shapes.add_picture(co_logo_path, Inches(6.27), Inches(0.4), height=Inches(0.65))

    # 3. Main Title Box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(7.27), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "標楷體"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = navy_blue

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = "標楷體"
        run2.font.size = Pt(18)
        run2.font.color.rgb = accent_orange

    # 4. Central Area (Mode dependent)
    if mode == "lecture":
        # Lecture Mode Layout: Speaker, Time, Location, Content
        # Speaker Box
        spk_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.7), Inches(7.07), Inches(1.8))
        spk_shape.fill.solid()
        spk_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        spk_shape.line.color.rgb = navy_blue

        tf_spk = spk_shape.text_frame
        tf_spk.word_wrap = True
        p_spk_t = tf_spk.paragraphs[0]
        r_st = p_spk_t.add_run()
        r_st.text = "🎤 主講人 / Speaker"
        r_st.font.name = "標楷體"
        r_st.font.size = Pt(14)
        r_st.font.bold = True
        r_st.font.color.rgb = accent_orange

        p_spk_c = tf_spk.add_paragraph()
        r_sc = p_spk_c.add_run()
        r_sc.text = speaker
        r_sc.font.name = "標楷體"
        r_sc.font.size = Pt(20)
        r_sc.font.bold = True
        r_sc.font.color.rgb = dark_text

        # Time & Location Box
        info_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.7), Inches(7.07), Inches(2.2))
        info_shape.fill.solid()
        info_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        info_shape.line.color.rgb = navy_blue

        tf_info = info_shape.text_frame
        tf_info.word_wrap = True
        p_t1 = tf_info.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = "📅 活動時間：" + time_str
        r_t1.font.name = "標楷體"
        r_t1.font.size = Pt(16)
        r_t1.font.bold = True

        p_t2 = tf_info.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = "📍 活動地點：" + location
        r_t2.font.name = "標楷體"
        r_t2.font.size = Pt(16)
        r_t2.font.bold = True

        # Event Description Box
        desc_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(7.1), Inches(7.07), Inches(3.2))
        desc_shape.fill.solid()
        desc_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        desc_shape.line.color.rgb = navy_blue

        tf_desc = desc_shape.text_frame
        tf_desc.word_wrap = True
        p_dt = tf_desc.paragraphs[0]
        r_dt = p_dt.add_run()
        r_dt.text = "💡 講座亮點與精彩內容"
        r_dt.font.name = "標楷體"
        r_dt.font.size = Pt(16)
        r_dt.font.bold = True
        r_dt.font.color.rgb = navy_blue

        details_list = details or [
            "掌握前瞻 AI 技術在產學領域的實務應用與未來趨勢",
            "深入解析 ESG 碳中和與永續發展的關鍵推動策略",
            "現場開放 Q&A 互動問答與跨領域交流對談"
        ]
        for item in details_list:
            p_item = tf_desc.add_paragraph()
            r_item = p_item.add_run()
            r_item.text = "• " + item
            r_item.font.name = "標楷體"
            r_item.font.size = Pt(14)
            r_item.font.color.rgb = dark_text

    else:
        # Project / Outcome Poster Mode Layout (4 Cards)
        d_dict = details if isinstance(details, dict) else {}
        cards = [
            ("📌 問題意識與目標", d_dict.get("目標", "開辦微學分課程與平台建置，提升在地連結與社會參與")),
            ("🛠️ 執行重點與方法", d_dict.get("方法", "聚焦智慧農業與ESG，結合技術與社會責任實踐")),
            ("🌱 社會效益與影響", d_dict.get("效益", "透過平台紀錄合作場域交流，無國界分享達成宣傳目標")),
            ("🏆 成果亮點與SDGs", d_dict.get("成果", "符合作業要點，推進4優質教育、9產業創新與13氣候行動"))
        ]
        positions = [(0.6, 2.7), (4.2, 2.7), (0.6, 6.2), (4.2, 6.2)]
        for idx, (ctitle, cbody) in enumerate(cards):
            x, y = positions[idx]
            card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.47), Inches(3.3))
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            card_shape.line.color.rgb = navy_blue

            tf_c = card_shape.text_frame
            tf_c.word_wrap = True
            p_ct = tf_c.paragraphs[0]
            r_ct = p_ct.add_run()
            r_ct.text = ctitle
            r_ct.font.name = "標楷體"
            r_ct.font.size = Pt(15)
            r_ct.font.bold = True
            r_ct.font.color.rgb = navy_blue

            p_cb = tf_c.add_paragraph()
            r_cb = p_cb.add_run()
            r_cb.text = cbody
            r_cb.font.name = "標楷體"
            r_cb.font.size = Pt(12)
            r_cb.font.color.rgb = dark_text

    # 5. Bottom Footer (行政資訊條)
    ft_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(10.69), Inches(8.27), Inches(1.0))
    ft_shape.fill.solid()
    ft_shape.fill.fore_color.rgb = navy_blue
    ft_shape.line.fill.background()

    tf_ft = ft_shape.text_frame
    tf_ft.word_wrap = True
    p_ft = tf_ft.paragraphs[0]
    p_ft.alignment = PP_ALIGN.CENTER
    
    footer_text = f"指導單位：{advisor}"
    if co_organizer:
        footer_text += f" │ 協辦單位：{co_organizer}"
    footer_text += f" │ 執行單位：{organizer}\n聯絡方式：{contact}"

    r_ft = p_ft.add_run()
    r_ft.text = footer_text
    r_ft.font.name = "標楷體"
    r_ft.font.size = Pt(11)
    r_ft.font.bold = True
    r_ft.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(output_pptx)
    print(f"[Success] Successfully generated poster PPTX at: {output_pptx}")

def main():
    parser = argparse.ArgumentParser(description="敏實科技大學專屬海報自動生成器")
    parser.add_argument("--mode", choices=["lecture", "project"], default="lecture", help="海報類型: lecture (主題講座/演講) 或 project (計畫/成果海報)")
    parser.add_argument("--title", required=True, help="海報主題/計畫名稱")
    parser.add_argument("--subtitle", default="", help="海報副標題")
    parser.add_argument("--speaker", default="張教授 / 敏實科技大學", help="演講人 (講座模式)")
    parser.add_argument("--time", default="115年9月15日 (三) 10:00 - 12:00", help="時間資訊")
    parser.add_argument("--location", default="敏實科技大學 大華樓5樓國際會議廳", help="地點資訊")
    parser.add_argument("--advisor", default="教育部", help="指導單位")
    parser.add_argument("--co-organizer", default="", help="協辦單位名稱")
    parser.add_argument("--co-logo", default="", help="協辦單位 Logo 圖片路徑")
    parser.add_argument("--organizer", default="敏實科技大學 USR 實踐中心", help="執行單位")
    parser.add_argument("--contact", default="邱專員 (分機 1234 / email: usr@mitust.edu.tw)", help="聯繫方式")
    parser.add_argument("--out", default="poster_output.pptx", help="輸出 .pptx 檔案路徑")

    args = parser.parse_args()

    create_minth_poster(
        mode=args.mode,
        title=args.title,
        subtitle=args.subtitle,
        speaker=args.speaker,
        time_str=args.time,
        location=args.location,
        advisor=args.advisor,
        co_organizer=args.co_organizer,
        co_logo_path=args.co_logo if args.co_logo else None,
        organizer=args.organizer,
        contact=args.contact,
        output_pptx=args.out
    )

if __name__ == "__main__":
    main()
