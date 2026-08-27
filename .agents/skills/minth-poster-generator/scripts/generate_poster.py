#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
敏實科技大學專屬海報自動生成器 (generate_poster.py)
支援 5 款全新 Canva 模板風格：
1. eco: 【永續綠能】生態風 (Eco-Sustainability)
2. tech: 【前瞻科技】數據風 (Future Tech)
3. agenda: 【實作工作坊】日程表風 (Workshop Agenda)
4. keynote: 【焦點人物】大師開講風 (Keynote Speaker)
5. grid: 【成果展板】展板風 (Showcase Grid)

輸出格式：
- .pptx 檔案：可在 Canva 中直接匯入並保持向量圖層可編輯性
"""

import sys
import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT_LOGO_PATH = os.path.join(SKILL_DIR, "resources", "minth_logo.png")

def create_minth_poster(
    template="eco",
    title="AI與永續發展主題講座",
    subtitle="",
    speaker="張教授 / 敏實科技大學",
    time_str="115年9月15日 (三) 10:00 - 12:00",
    location="敏實科技大學 大華樓5樓國際會議廳",
    advisor="教育部",
    co_organizer="",
    co_logo_path=None,
    organizer="敏實科技大學 USR 實踐中心",
    contact="邱專員 (分機 1234 / email: usr@mitust.edu.tw)",
    details=None,
    output_pptx="poster_output.pptx"
):
    prs = Presentation()
    # A4 Portrait: 8.27 x 11.69 inches
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ================= 模板風格定義 =================
    if template == "eco" or template == "lecture":
        bg_color = RGBColor(0xF0, 0xF8, 0xE2)  # Soft Light Green
        primary_color = RGBColor(0x2E, 0x7D, 0x32) # Forest Green
        secondary_color = RGBColor(0x5D, 0x40, 0x37) # Earth Brown
        font_family = "標楷體"
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    elif template == "tech":
        bg_color = RGBColor(0xF0, 0xF4, 0xF8) # Light Blue Grey
        primary_color = RGBColor(0x00, 0x4B, 0x87) # Navy Blue
        secondary_color = RGBColor(0x00, 0x97, 0xA7) # Cyan
        font_family = "微軟正黑體"
        shape_type = MSO_SHAPE.RECTANGLE
    elif template == "agenda":
        bg_color = RGBColor(0xFF, 0xFF, 0xFF) # White
        primary_color = RGBColor(0x37, 0x47, 0x4F) # Blue Grey
        secondary_color = RGBColor(0x54, 0x6E, 0x7A)
        font_family = "微軟正黑體"
        shape_type = MSO_SHAPE.RECTANGLE
    elif template == "keynote":
        bg_color = RGBColor(0xFA, 0xFA, 0xFA)
        primary_color = RGBColor(0x88, 0x0E, 0x4F) # Burgundy
        secondary_color = RGBColor(0xFF, 0xA0, 0x00) # Gold
        font_family = "標楷體"
        shape_type = MSO_SHAPE.HEXAGON
    elif template == "grid" or template == "project":
        bg_color = RGBColor(0xFF, 0xFF, 0xFF)
        primary_color = RGBColor(0x42, 0x42, 0x42) # Dark Grey
        secondary_color = RGBColor(0x75, 0x75, 0x75)
        font_family = "微軟正黑體"
        shape_type = MSO_SHAPE.RECTANGLE
    else:
        bg_color = RGBColor(0xF4, 0xF8, 0xFB)
        primary_color = RGBColor(0x00, 0x4B, 0x87)
        secondary_color = RGBColor(0xE6, 0x51, 0x00)
        font_family = "標楷體"
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

    # 1. Background Fill
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8.27), Inches(11.69))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.fill.background()

    # 2. Header Logos
    if os.path.exists(DEFAULT_LOGO_PATH):
        slide.shapes.add_picture(DEFAULT_LOGO_PATH, Inches(0.4), Inches(0.4), height=Inches(0.65))
    if co_logo_path and os.path.exists(co_logo_path):
        slide.shapes.add_picture(co_logo_path, Inches(6.27), Inches(0.4), height=Inches(0.65))

    # 3. Main Title Box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(7.27), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = font_family
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = primary_color

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = font_family
        run2.font.size = Pt(18)
        run2.font.color.rgb = secondary_color

    # 4. Central Area Logic
    if template in ["eco", "tech", "lecture"]:
        # Standard Lecture Layout
        spk_shape = slide.shapes.add_shape(shape_type, Inches(0.6), Inches(2.7), Inches(7.07), Inches(1.8))
        spk_shape.fill.solid()
        spk_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        spk_shape.line.color.rgb = primary_color

        tf_spk = spk_shape.text_frame
        tf_spk.word_wrap = True
        p_spk_t = tf_spk.paragraphs[0]
        r_st = p_spk_t.add_run()
        r_st.text = "🎤 主講人 / Speaker"
        r_st.font.name = font_family
        r_st.font.size = Pt(14)
        r_st.font.bold = True
        r_st.font.color.rgb = secondary_color

        p_spk_c = tf_spk.add_paragraph()
        r_sc = p_spk_c.add_run()
        r_sc.text = speaker
        r_sc.font.name = font_family
        r_sc.font.size = Pt(20)
        r_sc.font.bold = True
        r_sc.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

        info_shape = slide.shapes.add_shape(shape_type, Inches(0.6), Inches(4.7), Inches(7.07), Inches(2.2))
        info_shape.fill.solid()
        info_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        info_shape.line.color.rgb = primary_color

        tf_info = info_shape.text_frame
        tf_info.word_wrap = True
        p_t1 = tf_info.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = "📅 活動時間：" + time_str
        r_t1.font.name = font_family
        r_t1.font.size = Pt(16)
        r_t1.font.bold = True

        p_t2 = tf_info.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = "📍 活動地點：" + location
        r_t2.font.name = font_family
        r_t2.font.size = Pt(16)
        r_t2.font.bold = True

        desc_shape = slide.shapes.add_shape(shape_type, Inches(0.6), Inches(7.1), Inches(7.07), Inches(3.2))
        desc_shape.fill.solid()
        desc_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        desc_shape.line.color.rgb = primary_color

        tf_desc = desc_shape.text_frame
        tf_desc.word_wrap = True
        p_dt = tf_desc.paragraphs[0]
        r_dt = p_dt.add_run()
        r_dt.text = "💡 講座亮點與精彩內容"
        r_dt.font.name = font_family
        r_dt.font.size = Pt(16)
        r_dt.font.bold = True
        r_dt.font.color.rgb = primary_color

    elif template == "agenda":
        # Workshop Agenda Table Layout
        # Add basic info above table
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(7.27), Inches(1.0))
        tf_info = info_box.text_frame
        p_info = tf_info.paragraphs[0]
        p_info.alignment = PP_ALIGN.CENTER
        r_info = p_info.add_run()
        r_info.text = f"時間：{time_str} | 地點：{location}"
        r_info.font.name = font_family
        r_info.font.size = Pt(14)
        
        rows, cols = 5, 3
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3.5), Inches(7.27), Inches(5.0))
        table = table_shape.table
        
        headers = ["時間", "課程單元", "講師"]
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = primary_color
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_family
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        sample_agenda = [
            ("09:30 - 10:00", "報到與開場", "主持人"),
            ("10:00 - 12:00", "主題演講上半場", speaker),
            ("12:00 - 13:00", "午餐休息", "-"),
            ("13:00 - 16:00", "實作工作坊", speaker)
        ]
        
        for r_idx, row_data in enumerate(sample_agenda):
            for c_idx, text in enumerate(row_data):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = text
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font_family
                        run.font.size = Pt(12)

    elif template == "keynote":
        # Avatar placeholder
        avatar_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.63), Inches(2.5), Inches(3.0), Inches(3.0))
        avatar_shape.fill.solid()
        avatar_shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        avatar_shape.line.fill.background()
        tf_av = avatar_shape.text_frame
        p_av = tf_av.paragraphs[0]
        p_av.alignment = PP_ALIGN.CENTER
        r_av = p_av.add_run()
        r_av.text = "(照片置換區)"
        r_av.font.name = font_family

        spk_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.13), Inches(5.8), Inches(6.0), Inches(1.5))
        spk_banner.fill.solid()
        spk_banner.fill.fore_color.rgb = primary_color
        tf_bn = spk_banner.text_frame
        p_bn = tf_bn.paragraphs[0]
        p_bn.alignment = PP_ALIGN.CENTER
        r_bn = p_bn.add_run()
        r_bn.text = speaker
        r_bn.font.name = font_family
        r_bn.font.size = Pt(24)
        r_bn.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r_bn.font.bold = True
        
        info_box = slide.shapes.add_textbox(Inches(1.0), Inches(7.5), Inches(6.27), Inches(2.5))
        tf_info = info_box.text_frame
        p_i = tf_info.paragraphs[0]
        p_i.alignment = PP_ALIGN.CENTER
        r_i = p_i.add_run()
        r_i.text = f"時間：{time_str}\n地點：{location}"
        r_i.font.name = font_family
        r_i.font.size = Pt(18)
        r_i.font.color.rgb = primary_color

    elif template in ["grid", "project"]:
        # 4 Quadrants Grid
        cards = [
            ("📌 模塊一", "專題簡介與目標"),
            ("🛠️ 模塊二", "執行重點與方法"),
            ("🌱 模塊三", "數據與效益"),
            ("🏆 模塊四", "成果亮點與SDGs")
        ]
        positions = [(0.6, 2.7), (4.2, 2.7), (0.6, 6.2), (4.2, 6.2)]
        for idx, (ctitle, cbody) in enumerate(cards):
            x, y = positions[idx]
            card_shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(3.47), Inches(3.3))
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            card_shape.line.color.rgb = secondary_color
            card_shape.line.width = Pt(2)

            tf_c = card_shape.text_frame
            tf_c.word_wrap = True
            p_ct = tf_c.paragraphs[0]
            r_ct = p_ct.add_run()
            r_ct.text = ctitle
            r_ct.font.name = font_family
            r_ct.font.size = Pt(15)
            r_ct.font.bold = True
            r_ct.font.color.rgb = primary_color

            p_cb = tf_c.add_paragraph()
            r_cb = p_cb.add_run()
            r_cb.text = cbody
            r_cb.font.name = font_family
            r_cb.font.size = Pt(12)

    # 5. Bottom Footer (行政資訊條)
    ft_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(10.69), Inches(8.27), Inches(1.0))
    ft_shape.fill.solid()
    ft_shape.fill.fore_color.rgb = primary_color
    ft_shape.line.fill.background()

    tf_ft = ft_shape.text_frame
    tf_ft.word_wrap = True
    p_ft = tf_ft.paragraphs[0]
    p_ft.alignment = PP_ALIGN.CENTER
    
    footer_text = f"指導單位：{advisor}"
    if co_organizer:
        footer_text += f" │ 協辦單位：{co_organizer}"
    footer_text += f" │ 執行單位：{organizer}\n聯絡方式：{contact}"

    r_ft = p_ft.add_run()
    r_ft.text = footer_text
    r_ft.font.name = font_family
    r_ft.font.size = Pt(11)
    r_ft.font.bold = True
    r_ft.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(output_pptx)
    print(f"[Success] Successfully generated {template} poster PPTX at: {output_pptx}")

def main():
    parser = argparse.ArgumentParser(description="敏實科技大學專屬海報自動生成器")
    parser.add_argument("--template", choices=["eco", "tech", "agenda", "keynote", "grid", "lecture", "project"], default="eco", help="海報模板風格")
    # 兼容舊參數
    parser.add_argument("--mode", help=argparse.SUPPRESS)
    
    parser.add_argument("--title", required=True, help="海報主題/計畫名稱")
    parser.add_argument("--subtitle", default="", help="海報副標題")
    parser.add_argument("--speaker", default="張教授 / 敏實科技大學", help="演講人")
    parser.add_argument("--time", default="115年9月15日 (三) 10:00 - 12:00", help="時間資訊")
    parser.add_argument("--location", default="敏實科技大學 大華樓5樓國際會議廳", help="地點資訊")
    parser.add_argument("--advisor", default="教育部", help="指導單位")
    parser.add_argument("--co-organizer", default="", help="協辦單位名稱")
    parser.add_argument("--co-logo", default="", help="協辦單位 Logo 圖片路徑")
    parser.add_argument("--organizer", default="敏實科技大學 USR 實踐中心", help="執行單位")
    parser.add_argument("--contact", default="邱專員 (分機 1234 / email: usr@mitust.edu.tw)", help="聯繫方式")
    parser.add_argument("--out", default="poster_output.pptx", help="輸出 .pptx 檔案路徑")

    args = parser.parse_args()

    # 相容舊指令
    active_template = args.template
    if args.mode:
        active_template = args.mode

    create_minth_poster(
        template=active_template,
        title=args.title,
        subtitle=args.subtitle,
        speaker=args.speaker,
        time_str=args.time,
        location=args.location,
        advisor=args.advisor,
        co_organizer=args.co_organizer,
        co_logo_path=args.co_logo if args.co_logo else None,
        organizer=args.organizer,
        contact=args.contact,
        output_pptx=args.out
    )

if __name__ == "__main__":
    main()
